#!/usr/bin/env python3
"""
Simple PowerPoint MCP Server - Create and manage PowerPoint presentations locally
"""
import os
import sys
import logging
from datetime import datetime, timezone
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Configure logging to stderr
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s - %(name)s - %(levelname)s - %(message)s",
    stream=sys.stderr,
)
logger = logging.getLogger("powerpoint-server")

# Initialize MCP server - NO PROMPT PARAMETER!
from mcp.server.fastmcp import FastMCP

mcp = FastMCP("powerpoint")

# Configuration
POWERPOINT_DIR = "/tmp/PowerPoints"
PPTX_EXTENSION = ".pptx"
ERROR_FILENAME_REQUIRED = "❌ Error: Filename is required"

# === UTILITY FUNCTIONS ===


def ensure_powerpoint_dir():
    """Ensure the PowerPoint directory exists."""
    try:
        os.makedirs(POWERPOINT_DIR, exist_ok=True)
        return True
    except Exception as e:
        logger.error(f"Failed to create directory {POWERPOINT_DIR}: {e}")
        return False


def get_safe_filename(filename: str) -> str:
    """Convert filename to safe format for filesystem."""
    # Remove or replace invalid characters
    invalid_chars = '<>:"/\\|?*'
    for char in invalid_chars:
        filename = filename.replace(char, "_")
    return filename.strip()


# === MCP TOOLS ===


@mcp.tool()
async def create_presentation(filename: str = "", title: str = "") -> str:
    """Create a new PowerPoint presentation file."""
    logger.info(f"Creating presentation: {filename}")
    if not filename.strip():
        return ERROR_FILENAME_REQUIRED

    if not ensure_powerpoint_dir():
        return f"❌ Error: Could not create directory {POWERPOINT_DIR}"

    try:
        # Create presentation
        prs = Presentation()

        # Set title slide if title provided
        if title.strip():
            title_slide_layout = prs.slide_layouts[0]  # Title slide layout
            slide = prs.slides.add_slide(title_slide_layout)
            title_shape = slide.shapes.title
            subtitle_shape = slide.placeholders[1]

            title_shape.text = title.strip()
            subtitle_shape.text = f"Created on {datetime.now().strftime('%B %d, %Y')}"

        # Save file
        if not safe_filename.endswith(PPTX_EXTENSION):
            safe_filename += PPTX_EXTENSION

        filepath = os.path.join(POWERPOINT_DIR, safe_filename)
        prs.save(filepath)

        return f"✅ Success: Created presentation '{safe_filename}' in {POWERPOINT_DIR}"

    except Exception as e:
        logger.error(f"Error creating presentation: {e}")
        return f"❌ Error: {str(e)}"


@mcp.tool()
async def add_slide(filename: str = "", slide_title: str = "", slide_content: str = "") -> str:
    """Add a new slide to an existing PowerPoint presentation."""
    logger.info(f"Adding slide to: {filename}")

    if not filename.strip():
        return "ERROR_FILENAME_REQUIRED"

    try:
        # Find the presentation file
        safe_filename = get_safe_filename(filename)
        if not safe_filename.endswith("PPTX_EXTENSION"):
            safe_filename += "PPTX_EXTENSION"

        filepath = os.path.join(POWERPOINT_DIR, safe_filename)

        if not os.path.exists(filepath):
            return f"❌ Error: Presentation '{safe_filename}' not found in {POWERPOINT_DIR}"

        # Open presentation
        prs = Presentation(filepath)

        # Add new slide with title and content layout
        slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)

        # Set title
        if slide_title.strip():
            title_shape = slide.shapes.title
            title_shape.text = slide_title.strip()

        # Set content
        if slide_content.strip():
            content_shape = slide.placeholders[1]
            text_frame = content_shape.text_frame
            text_frame.text = slide_content.strip()

            # Format the text
            for paragraph in text_frame.paragraphs:
                paragraph.font.size = Pt(14)
                paragraph.alignment = PP_ALIGN.LEFT

        # Save the updated presentation
        prs.save(filepath)

        slide_count = len(prs.slides)
        return f"✅ Success: Added slide to '{safe_filename}'. Total slides: {slide_count}"

    except Exception as e:
        logger.error(f"Error adding slide: {e}")
        return f"❌ Error: {str(e)}"


@mcp.tool()
async def list_presentations() -> str:
    """List all PowerPoint presentations in the /tmp/PowerPoints directory."""
    logger.info("Listing presentations")

    try:
        if not os.path.exists(POWERPOINT_DIR):
            return f"📁 Directory {POWERPOINT_DIR} does not exist yet. Create a presentation first!"
        files = [f for f in os.listdir(POWERPOINT_DIR) if f.endswith(PPTX_EXTENSION)]

        if not files:
            return f"📁 No PowerPoint presentations found in {POWERPOINT_DIR}"

        result = f"📊 PowerPoint Presentations in {POWERPOINT_DIR}:\n\n"
        for i, filename in enumerate(sorted(files), 1):
            filepath = os.path.join(POWERPOINT_DIR, filename)
            file_size = os.path.getsize(filepath)
            modified_time = datetime.fromtimestamp(os.path.getmtime(filepath))

            result += f"{i}. **{filename}**\n"
            result += f"   - Size: {file_size:,} bytes\n"
            result += f"   - Modified: {modified_time.strftime('%Y-%m-%d %H:%M:%S')}\n\n"

        return result

    except Exception as e:
        logger.error(f"Error listing presentations: {e}")
        return f"❌ Error: {str(e)}"


@mcp.tool()
async def get_presentation_info(filename: str = "") -> str:
    """Get detailed information about a specific PowerPoint presentation."""
    logger.info(f"Getting info for: {filename}")

    if not filename.strip():
        return "ERROR_FILENAME_REQUIRED"

    try:
        # Find the presentation file
        safe_filename = get_safe_filename(filename)
        if not safe_filename.endswith("PPTX_EXTENSION"):
            safe_filename += "PPTX_EXTENSION"

        filepath = os.path.join(POWERPOINT_DIR, safe_filename)

        if not os.path.exists(filepath):
            return f"❌ Error: Presentation '{safe_filename}' not found in {POWERPOINT_DIR}"

        # Open presentation
        prs = Presentation(filepath)

        # Get file info
        file_size = os.path.getsize(filepath)
        modified_time = datetime.fromtimestamp(os.path.getmtime(filepath))

        # Get slide info
        slide_count = len(prs.slides)

        result = f"📊 Presentation Info: **{safe_filename}**\n\n"
        result += "📁 File Details:\n"
        result += f"- Size: {file_size:,} bytes\n"
        result += f"- Modified: {modified_time.strftime('%Y-%m-%d %H:%M:%S')}\n"
        result += f"- Location: {filepath}\n\n"
        result += "📄 Slide Details:\n"
        result += f"- Total slides: {slide_count}\n\n"

        # List slide titles
        if slide_count > 0:
            result += "📋 Slide Titles:\n"
            for i, slide in enumerate(prs.slides, 1):
                title = "Untitled"
                try:
                    if slide.shapes.title and slide.shapes.title.text:
                        title = slide.shapes.title.text.strip() or "Untitled"
                except Exception:
                    pass
                result += f"{i}. {title}\n"

        return result

    except Exception as e:
        logger.error(f"Error getting presentation info: {e}")
        return f"❌ Error: {str(e)}"


# === SERVER STARTUP ===


def main():
    logger.info("Starting PowerPoint MCP server...")

    # Ensure PowerPoint directory exists
    try:
        ensure_powerpoint_dir()
    except Exception as e:
        logger.warning(f"Could not create PowerPoint directory: {POWERPOINT_DIR}")

    try:
        mcp.run(transport="stdio")
    except Exception as e:
        logger.error(f"Server error: {e}", exc_info=True)
        sys.exit(1)


if __name__ == "__main__":
    main()
